import argparse
import pathlib
import json
import hashlib
import re
from datetime import datetime
from typing import List, Dict, Tuple, Optional
import frontmatter
from tqdm import tqdm

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from pdfminer.high_level import extract_text as pdf_extract_text
except ImportError:
    pdf_extract_text = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None


def extract_text_from_pdf(path: pathlib.Path) -> str:
    if pdf_extract_text is None:
        raise ImportError("pdfminer.six не установлен. Установите: pip install pdfminer.six")
    text = pdf_extract_text(str(path))
    if not text or not text.strip():
        raise ValueError(f"Не удалось извлечь текст из PDF: {path}")
    return text


def extract_text_from_docx(path: pathlib.Path) -> str:
    doc = docx.Document(str(path))
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)
    return "\n\n".join(paragraphs)


def extract_text_from_pptx(path: pathlib.Path) -> str:
    prs = Presentation(str(path))
    slides_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"# Слайд {slide_num}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        slides_text.append("\n\n".join(slide_text))
    return "\n\n---\n\n".join(slides_text)


def extract_text_from_html(path: pathlib.Path) -> str:
    html_content = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html_content, "html.parser")
    for script in soup(["script", "style"]):
        script.decompose()
    return soup.get_text(separator="\n\n", strip=True)


def extract_text_from_file(path: pathlib.Path) -> Tuple[str, Dict]:
    suffix = path.suffix.lower()
    metadata = {
        "source": str(path),
        "file_type": suffix[1:] if suffix else "unknown",
        "parsed_at": datetime.utcnow().isoformat() + "Z"
    }

    try:
        if suffix == ".md":
            text = path.read_text(encoding="utf-8")
        elif suffix == ".pdf":
            if pdf_extract_text is None:
                raise ImportError(f"pdfminer.six не установлен для обработки {path.name}. Установите: pip install pdfminer.six")
            text = extract_text_from_pdf(path)
            metadata["pages"] = "unknown"
        elif suffix in [".docx", ".doc"]:
            if docx is None:
                raise ImportError(f"python-docx не установлен для обработки {path.name}. Установите: pip install python-docx")
            text = extract_text_from_docx(path)
        elif suffix == ".pptx":
            if Presentation is None:
                raise ImportError(f"python-pptx не установлен для обработки {path.name}. Установите: pip install python-pptx")
            text = extract_text_from_pptx(path)
            if Presentation:
                prs = Presentation(str(path))
                metadata["slides"] = len(prs.slides)
        elif suffix in [".html", ".htm"]:
            if BeautifulSoup is None:
                raise ImportError(f"beautifulsoup4 не установлен для обработки {path.name}. Установите: pip install beautifulsoup4")
            text = extract_text_from_html(path)
        else:
            text = path.read_text(encoding="utf-8", errors="ignore")

        if not text or not text.strip():
            raise ValueError(f"Извлечен пустой текст из файла: {path.name}")

        return text, metadata
    except Exception as e:
        raise Exception(f"Ошибка при извлечении текста из {path.name}: {e}") from e


def to_markdown(path: pathlib.Path, out_dir: pathlib.Path) -> pathlib.Path:
    text, metadata = extract_text_from_file(path)

    # Создаем Markdown с frontmatter вручную
    md_lines = ["---"]
    for key, value in metadata.items():
        if isinstance(value, str):
            md_lines.append(f"{key}: {value}")
        else:
            md_lines.append(f"{key}: {json.dumps(value, ensure_ascii=False)}")
    md_lines.append("---")
    md_lines.append("")
    md_lines.append(text)
    md = "\n".join(md_lines)

    out_path = out_dir / (path.stem + ".md")
    out_path.write_text(md, encoding="utf-8")
    return out_path


def simple_tokenize(text: str) -> List[str]:
    return text.split()


def chunk_text_simple(text: str, chunk_size: int, overlap: int) -> List[str]:
    tokens = simple_tokenize(text)
    chunks = []
    i = 0
    while i < len(tokens):
        chunk = tokens[i:i+chunk_size]
        chunks.append(" ".join(chunk))
        i += chunk_size - overlap
    return chunks


def chunk_text_markdown(text: str, chunk_size: int, overlap: int, include_headers: bool = True) -> List[Tuple[str, str]]:
    header_pattern = r'^(#{1,3})\s+(.+)$'
    lines = text.split('\n')
    chunks = []
    current_chunk = []
    current_heading = ""
    current_tokens = 0
    for line in lines:
        header_match = re.match(header_pattern, line.strip())
        if header_match:
            if current_chunk and current_tokens >= chunk_size:
                chunk_text = '\n'.join(current_chunk)
                chunks.append((chunk_text, current_heading))

                if overlap > 0:
                    overlap_tokens = simple_tokenize(chunk_text)[-overlap:]
                    current_chunk = [' '.join(overlap_tokens)]
                    current_tokens = len(overlap_tokens)
                else:
                    current_chunk = []
                    current_tokens = 0

            current_heading = header_match.group(2)
            if include_headers:
                current_chunk.append(line)
                current_tokens += len(simple_tokenize(line))
        else:
            current_chunk.append(line)
            current_tokens += len(simple_tokenize(line))

            if current_tokens >= chunk_size:
                chunk_text = '\n'.join(current_chunk)
                chunks.append((chunk_text, current_heading))

                if overlap > 0:
                    overlap_tokens = simple_tokenize(chunk_text)[-overlap:]
                    current_chunk = [' '.join(overlap_tokens)]
                    current_tokens = len(overlap_tokens)
                else:
                    current_chunk = []
                    current_tokens = 0

    if current_chunk:
        chunk_text = '\n'.join(current_chunk)
        chunks.append((chunk_text, current_heading))
    return chunks


def chunk_text_recursive(text: str, chunk_size: int, overlap: int, separators: List[str] = None) -> List[str]:
    if separators is None:
        separators = ["\n\n", "\n", ". ", " ", ""]

    def _recursive_split(text: str, separators: List[str]) -> List[str]:
        if not separators:
            return [text]

        separator = separators[0]
        if separator == "":
            return chunk_text_simple(text, chunk_size, overlap)

        splits = text.split(separator)
        chunks = []
        current_chunk = []
        current_tokens = 0

        for split in splits:
            split_tokens = len(simple_tokenize(split))

            if current_tokens + split_tokens > chunk_size and current_chunk:
                chunk_text = separator.join(current_chunk)
                chunks.append(chunk_text)

                if overlap > 0:
                    overlap_tokens = simple_tokenize(chunk_text)[-overlap:]
                    current_chunk = [' '.join(overlap_tokens)]
                    current_tokens = len(overlap_tokens)
                else:
                    current_chunk = []
                    current_tokens = 0

            current_chunk.append(split)
            current_tokens += split_tokens

        if current_chunk:
            chunks.append(separator.join(current_chunk))

        result = []
        for chunk in chunks:
            if len(simple_tokenize(chunk)) > chunk_size:
                result.extend(_recursive_split(chunk, separators[1:]))
            else:
                result.append(chunk)

        return result

    return _recursive_split(text, separators)


def main(args):
    raw_dir = pathlib.Path(args.raw)
    md_out = pathlib.Path(args.md_out)
    md_out.mkdir(parents=True, exist_ok=True)
    chunks_out = pathlib.Path(args.chunks_out)
    chunks_out.mkdir(parents=True, exist_ok=True)

    if args.rebuild:
        import shutil
        if md_out.exists():
            shutil.rmtree(md_out)
        md_out.mkdir(parents=True, exist_ok=True)
        if chunks_out.exists():
            shutil.rmtree(chunks_out)
        chunks_out.mkdir(parents=True, exist_ok=True)

    files_to_process = [
        f for f in raw_dir.rglob("*")
        if f.is_file() and not f.name.startswith('.')
    ]

    if not files_to_process:
        print(f"Не найдено файлов для обработки в {raw_dir}")
        return

    print(f"Найдено файлов для обработки: {len(files_to_process)}")
    for f in files_to_process:
        print(f"  - {f.name} ({f.suffix})")
    print()

    all_chunks = []
    errors = []

    with tqdm(total=len(files_to_process), desc="Парсинг файлов", unit="файл") as pbar:
        for f in files_to_process:
            pbar.set_postfix(file=f.name[:30])
            try:
                # Извлечение текста
                text, metadata = extract_text_from_file(f)

                if not text or not text.strip():
                    error_msg = f"Пустой текст в файле {f.name}"
                    tqdm.write(f"⚠ {error_msg}")
                    errors.append((f.name, error_msg))
                    pbar.update(1)
                    continue

                # Сохранение в Markdown с frontmatter
                md_lines = ["---"]
                for key, value in metadata.items():
                    if isinstance(value, str):
                        md_lines.append(f"{key}: {value}")
                    else:
                        md_lines.append(f"{key}: {json.dumps(value, ensure_ascii=False)}")
                md_lines.append("---")
                md_lines.append("")
                md_lines.append(text)
                md = "\n".join(md_lines)

                md_file = md_out / (f.stem + ".md")
                md_file.write_text(md, encoding="utf-8")

                # Загрузка для проверки
                post = frontmatter.load(md_file)
                text = post.content
                metadata = post.metadata

                if args.splitter == "markdown":
                    chunk_data = chunk_text_markdown(
                        text,
                        args.chunk_size,
                        args.overlap,
                        include_headers=args.include_headers
                    )
                    for idx, (chunk_text, heading) in enumerate(chunk_data):
                        uid = hashlib.sha1((str(md_file) + str(idx)).encode()).hexdigest()
                        obj = {
                            "id": uid,
                            "text": chunk_text,
                            "file_path": str(metadata.get("source", "")),
                            "heading": heading,
                            "metadata": {
                                "source": metadata.get("source", ""),
                                "md_file": str(md_file),
                                "chunk_index": idx,
                                **{k: v for k, v in metadata.items() if k != "source"}
                            }
                        }
                        all_chunks.append(obj)
                elif args.splitter == "recursive":
                    chunks = chunk_text_recursive(text, args.chunk_size, args.overlap)
                    for idx, chunk_text in enumerate(chunks):
                        uid = hashlib.sha1((str(md_file) + str(idx)).encode()).hexdigest()
                        obj = {
                            "id": uid,
                            "text": chunk_text,
                            "file_path": str(metadata.get("source", "")),
                            "heading": "",
                            "metadata": {
                                "source": metadata.get("source", ""),
                                "md_file": str(md_file),
                                "chunk_index": idx,
                                **{k: v for k, v in metadata.items() if k != "source"}
                            }
                        }
                        all_chunks.append(obj)
                else:
                    chunks = chunk_text_simple(text, args.chunk_size, args.overlap)
                    for idx, chunk_text in enumerate(chunks):
                        uid = hashlib.sha1((str(md_file) + str(idx)).encode()).hexdigest()
                        obj = {
                            "id": uid,
                            "text": chunk_text,
                            "file_path": str(metadata.get("source", "")),
                            "heading": "",
                            "metadata": {
                                "source": metadata.get("source", ""),
                                "md_file": str(md_file),
                                "chunk_index": idx,
                                **{k: v for k, v in metadata.items() if k != "source"}
                            }
                        }
                        all_chunks.append(obj)
            except Exception as e:
                error_msg = f"Ошибка при обработке {f.name}: {str(e)}"
                tqdm.write(f"{error_msg}")
                errors.append((f.name, str(e)))
                import traceback
                tqdm.write(f"   Детали: {traceback.format_exc().split(chr(10))[-2]}")
                pbar.update(1)
                continue

            pbar.update(1)
    output_jsonl = chunks_out / "chunks.jsonl"
    with open(output_jsonl, "w", encoding="utf-8") as f:
        with tqdm(total=len(all_chunks), desc="Сохранение чанков", unit="чанк") as pbar:
            for chunk in all_chunks:
                f.write(json.dumps(chunk, ensure_ascii=False) + "\n")
                pbar.update(1)
    print(f"\n{'='*60}")
    print(f"Результаты обработки:")
    print(f"{'='*60}")
    print(f"Обработано файлов: {len(list(md_out.glob('*.md')))}")
    print(f"Создано чанков: {len(all_chunks)}")
    print(f"Чанки сохранены в: {output_jsonl}")
    if errors:
        print(f"\nОшибки при обработке ({len(errors)} файлов):")
        for filename, error in errors:
            print(f"  - {filename}: {error}")
    else:
        print(f"\n✓ Все файлы успешно обработаны!")


if __name__ == "__main__":
    p = argparse.ArgumentParser(description="Парсинг документов и разбиение на чанки")
    p.add_argument("--raw", default="data/raw", help="Директория с исходными файлами")
    p.add_argument("--md_out", default="data/md_parsed", help="Директория для Markdown файлов")
    p.add_argument("--chunks_out", default="data/chunks", help="Директория для чанков")
    p.add_argument("--chunk_size", type=int, default=500, help="Размер чанка в токенах (100-1000)")
    p.add_argument("--overlap", type=int, default=100, help="Перекрытие между чанками")
    p.add_argument(
        "--splitter",
        choices=["simple", "markdown", "recursive"],
        default="recursive",
        help="Стратегия разбиения: simple, markdown (по заголовкам), recursive"
    )
    p.add_argument(
        "--include_headers",
        action="store_true",
        default=True,
        help="Включать заголовки в текст чанка (для markdown splitter)"
    )
    p.add_argument(
        "--rebuild",
        action="store_true",
        help="Пересобрать индекс (очистить выходные директории)"
    )
    args = p.parse_args()
    main(args)
