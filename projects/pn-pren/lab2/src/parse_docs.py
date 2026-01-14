import os
import json
from pathlib import Path
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from bs4 import BeautifulSoup
import markdown


def parse_docx(file_path):
    doc = Document(file_path)
    content = []
    
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            level = int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1
            content.append(f"{'#' * level} {para.text}")
        else:
            if para.text.strip():
                content.append(para.text)
    
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append('| ' + ' | '.join(cells) + ' |')
        if rows:
            content.append('\n'.join(rows))
    
    metadata = {
        'source': file_path.name,
        'type': 'docx',
        'path': str(file_path)
    }
    
    return '\n\n'.join(content), metadata


def parse_pptx(file_path):
    prs = Presentation(file_path)
    content = []
    
    for i, slide in enumerate(prs.slides, 1):
        content.append(f"## Slide {i}")
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                content.append(shape.text)
        
        content.append('')
    
    metadata = {
        'source': file_path.name,
        'type': 'pptx',
        'path': str(file_path),
        'slides': len(prs.slides)
    }
    
    return '\n\n'.join(content), metadata


def parse_pdf(file_path):
    reader = PdfReader(file_path)
    content = []
    
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text.strip():
            content.append(f"## Page {i}\n\n{text}")
    
    metadata = {
        'source': file_path.name,
        'type': 'pdf',
        'path': str(file_path),
        'pages': len(reader.pages)
    }
    
    return '\n\n'.join(content), metadata


def parse_html(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
    
    for script in soup(['script', 'style']):
        script.decompose()
    
    text = soup.get_text()
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    
    metadata = {
        'source': file_path.name,
        'type': 'html',
        'path': str(file_path)
    }
    
    return '\n\n'.join(lines), metadata


def parse_markdown(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    metadata = {
        'source': file_path.name,
        'type': 'markdown',
        'path': str(file_path)
    }
    
    return content, metadata


def parse_document(file_path):
    file_path = Path(file_path)
    ext = file_path.suffix.lower()
    
    parsers = {
        '.docx': parse_docx,
        '.pptx': parse_pptx,
        '.pdf': parse_pdf,
        '.html': parse_html,
        '.htm': parse_html,
        '.md': parse_markdown
    }
    
    if ext not in parsers:
        raise ValueError(f"Unsupported format: {ext}")
    
    return parsers[ext](file_path)


def process_directory(input_dir, output_dir):
    input_dir = Path(input_dir)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    results = []
    
    for file_path in input_dir.rglob('*'):
        if file_path.is_file() and file_path.suffix.lower() in ['.docx', '.pptx', '.pdf', '.html', '.htm', '.md']:
            try:
                print(f"Processing {file_path.name}...")
                content, metadata = parse_document(file_path)
                
                output_file = output_dir / f"{file_path.stem}.md"
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write(content)
                
                metadata['output'] = str(output_file)
                results.append(metadata)
                
            except Exception as e:
                print(f"Error processing {file_path.name}: {e}")
    
    metadata_file = output_dir / 'metadata.json'
    with open(metadata_file, 'w', encoding='utf-8') as f:
        json.dump(results, f, indent=2)
    
    print(f"\nProcessed {len(results)} documents")
    return results


if __name__ == '__main__':
    import argparse
    
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', default='data/docs', help='Input directory')
    parser.add_argument('--output', default='output/markdown', help='Output directory')
    args = parser.parse_args()
    
    process_directory(args.input, args.output)
