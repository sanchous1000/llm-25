import os
from pathlib import Path
from typing import Dict, List, Optional
from datetime import datetime
import json

import pypdf
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import markdownify


class DocumentParser:
    def parse(self, file_path: Path) -> Dict:
        raise NotImplementedError
    
    def _extract_metadata(self, file_path: Path) -> Dict:
        stat = file_path.stat()
        return {
            "source": str(file_path),
            "filename": file_path.name,
            "path": str(file_path.relative_to(file_path.parents[2])),
            "size": stat.st_size,
            "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
        }


class PDFParser(DocumentParser):
    def parse(self, file_path: Path) -> Dict:
        metadata = self._extract_metadata(file_path)
        pages = []
        
        with open(file_path, "rb") as f:
            pdf_reader = pypdf.PdfReader(f)
            metadata["total_pages"] = len(pdf_reader.pages)
            metadata["author"] = pdf_reader.metadata.get("/Author", "")
            metadata["title"] = pdf_reader.metadata.get("/Title", "")
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                pages.append({
                    "page": page_num,
                    "text": text,
                    "metadata": {
                        **metadata,
                        "page": page_num,
                    }
                })
        
        return {
            "metadata": metadata,
            "pages": pages,
            "format": "pdf"
        }


class DOCXParser(DocumentParser):
    def parse(self, file_path: Path) -> Dict:
        metadata = self._extract_metadata(file_path)
        sections = []
        
        doc = Document(file_path)
        metadata["author"] = doc.core_properties.author or ""
        metadata["title"] = doc.core_properties.title or ""
        
        current_section = {"heading": "", "text": "", "paragraphs": []}
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            if para.style.name.startswith("Heading"):
                if current_section["text"]:
                    sections.append(current_section)
                current_section = {
                    "heading": text,
                    "text": text + "\n",
                    "paragraphs": [text],
                    "metadata": {
                        **metadata,
                        "section": len(sections) + 1,
                    }
                }
            else:
                current_section["text"] += text + "\n"
                current_section["paragraphs"].append(text)
        
        if current_section["text"]:
            sections.append(current_section)
        
        metadata["total_sections"] = len(sections)
        
        return {
            "metadata": metadata,
            "sections": sections,
            "format": "docx"
        }


class PPTXParser(DocumentParser):
    def parse(self, file_path: Path) -> Dict:
        metadata = self._extract_metadata(file_path)
        slides = []
        
        prs = Presentation(file_path)
        metadata["total_slides"] = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            text = "\n".join(slide_text)
            slides.append({
                "slide": slide_num,
                "text": text,
                "metadata": {
                    **metadata,
                    "slide": slide_num,
                }
            })
        
        return {
            "metadata": metadata,
            "slides": slides,
            "format": "pptx"
        }


class MarkdownParser(DocumentParser):
    def parse(self, file_path: Path) -> Dict:
        metadata = self._extract_metadata(file_path)
        
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        
        sections = []
        current_section = {"heading": "", "text": "", "metadata": {**metadata, "section": 1}}
        
        for line in content.split("\n"):
            if line.startswith("#"):
                if current_section["text"]:
                    sections.append(current_section)
                level = len(line) - len(line.lstrip("#"))
                heading = line.lstrip("#").strip()
                current_section = {
                    "heading": heading,
                    "text": line + "\n",
                    "metadata": {
                        **metadata,
                        "section": len(sections) + 1,
                        "heading_level": level,
                    }
                }
            else:
                current_section["text"] += line + "\n"
        
        if current_section["text"]:
            sections.append(current_section)
        
        metadata["total_sections"] = len(sections)
        
        return {
            "metadata": metadata,
            "sections": sections,
            "format": "markdown"
        }


class HTMLParser(DocumentParser):
    def parse(self, file_path: Path) -> Dict:
        metadata = self._extract_metadata(file_path)
        
        with open(file_path, "r", encoding="utf-8") as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, "html.parser")
        markdown_content = markdownify.markdownify(str(soup), heading_style="ATX")
        
        md_parser = MarkdownParser()
        sections = []
        current_section = {"heading": "", "text": "", "metadata": {**metadata, "section": 1}}
        
        for line in markdown_content.split("\n"):
            if line.startswith("#"):
                if current_section["text"]:
                    sections.append(current_section)
                level = len(line) - len(line.lstrip("#"))
                heading = line.lstrip("#").strip()
                current_section = {
                    "heading": heading,
                    "text": line + "\n",
                    "metadata": {
                        **metadata,
                        "section": len(sections) + 1,
                        "heading_level": level,
                    }
                }
            else:
                current_section["text"] += line + "\n"
        
        if current_section["text"]:
            sections.append(current_section)
        
        metadata["total_sections"] = len(sections)
        
        return {
            "metadata": metadata,
            "sections": sections,
            "format": "html"
        }


class ParserFactory:
    _parsers = {
        ".pdf": PDFParser,
        ".docx": DOCXParser,
        ".pptx": PPTXParser,
        ".md": MarkdownParser,
        ".markdown": MarkdownParser,
        ".html": HTMLParser,
        ".htm": HTMLParser,
    }
    
    @classmethod
    def get_parser(cls, file_path: Path) -> Optional[DocumentParser]:
        ext = file_path.suffix.lower()
        parser_class = cls._parsers.get(ext)
        if parser_class:
            return parser_class()
        return None
    
    @classmethod
    def parse_file(cls, file_path: Path, output_dir: Path) -> Optional[Path]:
        parser = cls.get_parser(file_path)
        if not parser:
            return None
        
        parsed = parser.parse(file_path)
        output_path = output_dir / f"{file_path.stem}.md"
        
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("---\n")
            f.write(f"source: {parsed['metadata'].get('source', '')}\n")
            f.write(f"filename: {parsed['metadata'].get('filename', '')}\n")
            f.write(f"format: {parsed['format']}\n")
            if "total_pages" in parsed["metadata"]:
                f.write(f"total_pages: {parsed['metadata']['total_pages']}\n")
            if "total_slides" in parsed["metadata"]:
                f.write(f"total_slides: {parsed['metadata']['total_slides']}\n")
            if "total_sections" in parsed["metadata"]:
                f.write(f"total_sections: {parsed['metadata']['total_sections']}\n")
            f.write("---\n\n")
            
            if "pages" in parsed:
                for page in parsed["pages"]:
                    f.write(f"## Page {page['page']}\n\n")
                    f.write(page["text"])
                    f.write("\n\n")
            elif "slides" in parsed:
                for slide in parsed["slides"]:
                    f.write(f"## Slide {slide['slide']}\n\n")
                    f.write(slide["text"])
                    f.write("\n\n")
            elif "sections" in parsed:
                for section in parsed["sections"]:
                    if section["heading"]:
                        level = section["metadata"].get("heading_level", 2)
                        f.write("#" * level + f" {section['heading']}\n\n")
                    f.write(section["text"])
                    f.write("\n\n")
        
        metadata_path = output_dir / f"{file_path.stem}_metadata.json"
        with open(metadata_path, "w", encoding="utf-8") as f:
            json.dump(parsed["metadata"], f, indent=2, ensure_ascii=False)
        
        return output_path

