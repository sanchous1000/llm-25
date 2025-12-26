import os
import yaml
import argparse
from pathlib import Path
from datetime import datetime
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, UnstructuredHTMLLoader
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. PPTX files will be skipped.")
try:
    import markdown
    from markdownify import markdownify
    HTML_AVAILABLE = True
except ImportError:
    HTML_AVAILABLE = False
    print("Warning: markdownify not installed. HTML files will use basic text extraction.")

# Define paths
BASE_DIR = Path(__file__).parent.parent
DATA_DIR = BASE_DIR / "data"
MARKDOWN_DIR = DATA_DIR / "markdown"

def get_file_metadata(file_path: Path) -> dict:
    """Extract file metadata including modification date."""
    metadata = {}
    try:
        stat = file_path.stat()
        # Get modification time
        mod_time = datetime.fromtimestamp(stat.st_mtime)
        metadata["date"] = mod_time.strftime("%Y-%m-%d")
        metadata["date_modified"] = mod_time.isoformat()
    except Exception as e:
        print(f"Warning: Could not extract date metadata for {file_path.name}: {e}")
    return metadata

def save_as_markdown(content, metadata, filename):
    """Saves content and metadata as a Markdown file with YAML frontmatter."""
    # Clean metadata values to ensure they are safe for YAML
    clean_metadata = {k: str(v) for k, v in metadata.items()}
    
    frontmatter = yaml.dump(clean_metadata, default_flow_style=False, sort_keys=False)
    md_content = f"---\n{frontmatter}---\n\n{content}"
    
    output_path = MARKDOWN_DIR / f"{filename}.md"
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(md_content)
    print(f"Saved: {output_path}")

def parse_documents():
    """Parses PDF and DOCX files in data/ and saves them as Markdown."""
    if not MARKDOWN_DIR.exists():
        MARKDOWN_DIR.mkdir(parents=True)

    for file_path in DATA_DIR.glob("*"):
        if file_path.is_dir() or file_path.parent.name == "markdown":
            continue

        print(f"Processing: {file_path.name}")
        
        try:
            if file_path.suffix.lower() == ".pdf":
                loader = PyPDFLoader(str(file_path))
                pages = loader.load()
                
                # Combine pages into one Markdown document
                full_text = ""
                for page in pages:
                    # Simple normalization: Add page markers
                    full_text += f"\n\n## Page {page.metadata.get('page', 'unknown') + 1}\n\n"
                    full_text += page.page_content
                
                metadata = {
                    "source": file_path.name,
                    "original_path": str(file_path),
                    "page_count": len(pages),
                    "type": "pdf"
                }
                # Add file metadata (date, etc.)
                metadata.update(get_file_metadata(file_path))
                save_as_markdown(full_text, metadata, file_path.stem)

            elif file_path.suffix.lower() == ".docx":
                loader = Docx2txtLoader(str(file_path))
                docs = loader.load()
                
                full_text = "\n\n".join([d.page_content for d in docs])
                metadata = {
                    "source": file_path.name,
                    "original_path": str(file_path),
                    "type": "docx"
                }
                # Add file metadata (date, etc.)
                metadata.update(get_file_metadata(file_path))
                save_as_markdown(full_text, metadata, file_path.stem)
            
            elif file_path.suffix.lower() == ".pptx":
                if not PPTX_AVAILABLE:
                    print(f"Skipping {file_path.name}: python-pptx not installed. Install with: pip install python-pptx")
                    continue
                
                # Parse PPTX file
                prs = Presentation(str(file_path))
                full_text = ""
                slide_count = 0
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_count += 1
                    full_text += f"\n\n## Slide {slide_num}\n\n"
                    
                    # Extract text from all shapes in the slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            full_text += shape.text.strip() + "\n\n"
                
                metadata = {
                    "source": file_path.name,
                    "original_path": str(file_path),
                    "slide_count": slide_count,
                    "type": "pptx"
                }
                # Add file metadata (date, etc.)
                metadata.update(get_file_metadata(file_path))
                save_as_markdown(full_text, metadata, file_path.stem)
                
            elif file_path.suffix.lower() in [".html", ".htm"]:
                if HTML_AVAILABLE:
                    # Use UnstructuredHTMLLoader for better HTML parsing
                    try:
                        loader = UnstructuredHTMLLoader(str(file_path))
                        docs = loader.load()
                        full_text = "\n\n".join([d.page_content for d in docs])
                    except:
                        # Fallback: read and convert HTML to markdown
                        with open(file_path, "r", encoding="utf-8") as f:
                            html_content = f.read()
                        full_text = markdownify(html_content, heading_style="ATX")
                else:
                    # Basic text extraction
                    loader = UnstructuredHTMLLoader(str(file_path))
                    docs = loader.load()
                    full_text = "\n\n".join([d.page_content for d in docs])
                
                metadata = {
                    "source": file_path.name,
                    "original_path": str(file_path),
                    "type": "html"
                }
                # Add file metadata (date, etc.)
                metadata.update(get_file_metadata(file_path))
                save_as_markdown(full_text, metadata, file_path.stem)
            
            elif file_path.suffix.lower() in [".md", ".markdown"]:
                # Just copy or wrap existing markdown
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                
                metadata = {
                    "source": file_path.name,
                    "original_path": str(file_path),
                    "type": "markdown"
                }
                # Add file metadata (date, etc.)
                metadata.update(get_file_metadata(file_path))
                save_as_markdown(content, metadata, file_path.stem)

        except Exception as e:
            print(f"Error processing {file_path.name}: {e}")

if __name__ == "__main__":
    parse_documents()
